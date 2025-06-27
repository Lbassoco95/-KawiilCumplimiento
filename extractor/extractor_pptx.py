from pptx import Presentation  # type: ignore
from typing import List

def extract_text_from_pptx(path: str) -> str:
    """
    Extrae el texto de un archivo .pptx, recorriendo cada diapositiva y extrayendo todo el texto visible.
    Args:
        path (str): Ruta al archivo .pptx.
    Returns:
        str: Texto extraído de todas las diapositivas del PowerPoint.
    """
    try:
        # Abrir la presentación
        prs = Presentation(path)
        all_text = []
        
        # Recorrer cada diapositiva
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Agregar título de la diapositiva
            slide_text.append(f"=== DIAPOSITIVA {slide_num} ===")
            
            # Extraer texto de cada forma en la diapositiva
            for shape in slide.shapes:
                # Verificar si la forma tiene texto
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Verificar si la forma es una tabla
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            slide_text.append(" | ".join(row_text))
            
            # Unir todo el texto de la diapositiva
            if len(slide_text) > 1:  # Si hay más que solo el título
                all_text.append("\n".join(slide_text))
        
        # Unir todo el texto de todas las diapositivas
        return "\n\n".join(all_text)
        
    except Exception as e:
        return f"Error al procesar el archivo PowerPoint: {str(e)}" 